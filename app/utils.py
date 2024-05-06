from threading import Lock
from whoosh.fields import Schema, TEXT, STORED, KEYWORD, ID
from whoosh.index import create_in, open_dir
from whoosh.query import *
from whoosh.analysis import StandardAnalyzer, KeywordAnalyzer
import os.path
from flask import current_app
from concurrent.futures import ThreadPoolExecutor
from collections import Counter
from spacy.lang.fr.stop_words import STOP_WORDS as fr_stop
from spacy.lang.en.stop_words import STOP_WORDS as en_stop
from spacy.util import compile_infix_regex
from spacy import load
from unidecode import unidecode
from pandas import read_excel, read_csv
from docx import Document
from PIL import Image, UnidentifiedImageError
import pytesseract
import fitz
from xlrd import open_workbook
from pptx import Presentation
from pdf2image import convert_from_path
import os
from chardet import detect
from app.extensions import db
from flask_login import current_user
from app.models.favoris import FAVORIS
from app.models.notification import NOTIFICATION
from app.models.fichier import FICHIER
from app.models.dossier import DOSSIER
import re
import zipfile
import os
import uuid

class SingletonMeta(type):
    """Singleton metaclass.

    Args:
        type (type): Metaclass type.

    Returns:
        type: Metaclass type.
    """
    _instances = {}
    _lock = Lock()
    
    def __call__(cls, *args, **kwargs):
        with cls._lock:
            if cls not in cls._instances:
                instance = super().__call__(*args, **kwargs)
                cls._instances[cls] = instance
        return cls._instances[cls]


class Whoosh(metaclass=SingletonMeta):
    def __init__(self):
        analyzer = StandardAnalyzer(stoplist=None)
        schema = Schema(title=TEXT(stored=True, analyzer=analyzer), content=TEXT(analyzer=analyzer), path=ID(analyzer=KeywordAnalyzer(), stored=True), tags=KEYWORD(commas=True, scorable=True, analyzer=analyzer), id=ID(stored=True))
        with current_app.app_context():
            if not os.path.exists(f'{current_app.root_path}/storage/index') :
                os.mkdir(f'{current_app.root_path}/storage/index')
                create_in(f'{current_app.root_path}/storage/index', schema)
            self.open_index = open_dir(f'{current_app.root_path}/storage/index')

    def add_document(self, title, content, path, tags, id):
        writer = self.open_index.writer()
        try:
            writer.add_document(title=title, content=content, path=path, tags=tags, id=id)
        finally:
            writer.commit()

    def delete_document(self, id):
        writer = self.open_index.writer()
        try:
            writer.delete_by_term("id", id)
        finally:
            writer.commit()

    def delete_documents(self, ids):
        writer = self.open_index.writer()
        try:
            for id in ids:
                writer.delete_by_term("id", id)
        finally:
            writer.commit()

    def search(self, query, path=None):
        if query.strip() == "":
            subquery = Every()
        else : 
            query = unidecode(query).lower()
            if query.startswith("&") or query.startswith("|"):
                query = "* " + query
            while re.search(r'([&|])\s*([&|]|$)', query):
                query = re.sub(r'([&|])\s*([&|]|$)', r'\1 * \2', query)
            or_conditions = [cond.strip() for cond in query.split("|")]
            conditions = [[cond.strip() for cond in condition.split('&')] for condition in or_conditions]
            if path is not None:
                path = path.strip().replace(" ", "")
                subquery = And([Term("path", path), Or([And([Or([Phrase("content", condition.split()), Term("tags", condition), Wildcard("title", "*"+condition.replace(" ", "_")+"*")]) for condition in condition_list]) for condition_list in conditions])])
            else:
                subquery = Or([And([Or([Phrase("content", condition.split()), Term("tags", condition), Wildcard("title", "*"+condition.replace(" ", "_")+"*")]) for condition in condition_list]) for condition_list in conditions])
        favoris = db.session.query(FAVORIS.c.id_Fichier).filter(FAVORIS.c.id_Utilisateur == current_user.id_Utilisateur).all()
        favoris_ids = [favori.id_Fichier for favori in favoris]
        with self.open_index.searcher() as searcher:
            results = searcher.search(subquery, limit=None)
            results_list = []
            for result in results:
                result_field = result.fields()
                result_field['extension'] = result_field['title'].split('.')[-1]
                if int(result_field['id']) in favoris_ids:
                    result_field['favori'] = True
                else :
                    result_field['favori'] = False
                results_list.append(result_field)
        return results_list


class NLPProcessor(metaclass=SingletonMeta):
    def __init__(self, batch_size=100000):
        self.batch_size = batch_size
        self.tokenizer_nlp = load(f"fr_application_model")
        self.lemmatizer_nlp = load("fr_core_news_sm")
        self.tokenizer_nlp.max_length = batch_size
        self.lemmatizer_nlp.max_length = batch_size
        default_infixes = list(self.tokenizer_nlp.Defaults.infixes)
        default_infixes.append("[A-Z][a-z0-9]+")
        infix_regex = compile_infix_regex(default_infixes)
        self.tokenizer_nlp.tokenizer.infix_finditer = infix_regex.finditer
        self.stop_words = set(fr_stop).union(en_stop)

    def clean(self, text):
        return [
            unidecode(token.text).lower()
            for token in self.tokenizer_nlp(text)
            if unidecode(token.text).lower() not in self.stop_words
            and len(token) >= 3
            and not token.is_stop
            and not token.is_punct
            and not token.is_space
            and not token.like_url
            and not token.like_email
            and not token.is_digit
            and not token.is_currency
        ]

    def lemmatize(self, text):
        batches = [
            text[i : i + self.batch_size] for i in range(0, len(text), self.batch_size)
        ]
        result = []
        with ThreadPoolExecutor() as executor:
            for batch in executor.map(self.process_lemmatize_batch, batches):
                result.append(" ".join(token.lemma_ for token in batch))
        return " ".join(result)

    def process_lemmatize_batch(self, batch):
        return self.lemmatizer_nlp(batch)

    def tokenize(self, text):
        text = self.lemmatize(text)
        batches = [
            text[i : i + self.batch_size] for i in range(0, len(text), self.batch_size)
        ]
        word_frequencies = {}
        with ThreadPoolExecutor() as executor:
            for batch in executor.map(self.process_tokenize_batch, batches):
                word_frequencies.update(batch)
        word_frequencies = sorted(
            word_frequencies.items(), key=lambda x: x[1], reverse=True
        )
        return [word for word, _ in word_frequencies]

    def process_tokenize_batch(self, batch):
        return Counter(self.clean(batch))


class FileReader(metaclass=SingletonMeta):
    def __init__(self):
        self.readers = {
            "csv": self.read_csv,
            #"doc": self.read_doc,
            "docx": self.read_docx,
            #"json": self.read_json,
            #"html": self.read_html,
            #"htm": self.read_html,
            #"odt": self.read_odt,
            "pdf": self.read_pdf,
            "txt": self.read_txt,
            "xlsx": self.read_xlsx,
            "xls": self.read_xls,
            "pptx": self.read_pptx,
            #"ppt": self.read_ppt,
            #"gif": self.read_ocr,
            "jpg": self.read_ocr,
            "jpeg": self.read_ocr,
            "png": self.read_ocr,
            "tiff": self.read_ocr,
            "tif": self.read_ocr,
        }

    def read(self, file_path, extension):
        extension = extension.lower()
        return self.readers.get(extension, lambda x: '')(file_path)
    
    def read_csv(self, file_path):
        if os.stat(file_path).st_size == 0:
            return None
        return read_csv(file_path).to_string(index=False)

    def read_docx(self, file_path):
        text = ""
        for paragraph in Document(file_path).paragraphs:
            text += paragraph.text
        return text
    
    def read_ocr(self, file_path):
        try:
            return pytesseract.image_to_string(Image.open(file_path), lang="fra")
        except UnidentifiedImageError:
            return None

    def read_pdf(self, file_path):
        file = fitz.open(file_path)
        text = ""
        for page in file:
            text += page.get_text()
            text.replace("\n", " ")
        cleaned_text = ''.join([i for i in text if i.isprintable() and not i.isspace()])
        if cleaned_text == "":
            images = convert_from_path(file_path)
            for image in images:
                text += pytesseract.image_to_string(image, lang="fra")
        return text

    def read_txt(self, file_path):
        with open(file_path, "r", encoding=detect(open(file_path, "rb").read())["encoding"]) as file:
            return file.read()

    def read_xlsx(self, file_path):
        df = read_excel(file_path, engine="openpyxl")
        text = ''
        for column in df.columns:
            text += df[column].astype(str).str.cat(sep=' ') + ' '
            text = text.replace('nan', '')
        return text

    def read_xls(self, file_path):
        file = open_workbook(file_path)
        text = ""
        for sheet in file.sheets():
            for row in range(sheet.nrows):
                for cell in sheet.row(row):
                    text += str(cell.value) + " "
        return text

    def read_pptx(self, file_path):
        text = ""
        for slide in Presentation(file_path).slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        return text

class FileDownloader(metaclass=SingletonMeta):
    def create_zip(self, file_ids):
        if not os.path.exists(f"{current_app.root_path}/storage/downloads"):
            os.makedirs(f"{current_app.root_path}/storage/downloads")
        zip_filename = f"{uuid.uuid4().hex}.zip"
        zip_path = f"{current_app.root_path}/storage/downloads/{zip_filename}"
        files = db.session.query(FICHIER).filter(FICHIER.id_Fichier.in_(file_ids)).all()
        files_by_folder = {}
        for file in files:
            if file.DOSSIER_.id_Dossier not in files_by_folder:
                files_by_folder[file.DOSSIER_.id_Dossier] = []
            files_by_folder[file.DOSSIER_.id_Dossier].append(file)
        with zipfile.ZipFile(zip_path, "w") as zip_file:
            for folder_id in files_by_folder:
                database_folder = db.session.query(DOSSIER).filter(DOSSIER.id_Dossier == folder_id).first()
                for file in files_by_folder[folder_id]:
                    file_path = os.path.join(current_app.root_path, 'storage', str(folder_id), f'{file.id_Fichier}.{file.extension_Fichier}')
                    zip_file.write(file_path, arcname=os.path.join(database_folder.nom_Dossier.replace('/', '-'), f'{file.nom_Fichier}.{file.extension_Fichier}'))
        return zip_path

def check_notitications():
    """Check if there is any notification in the database.

    Returns:
        bool: True if there is any notification, False otherwise.
    """
    return NOTIFICATION.query.all() != []

def get_total_file_count(folder):
    total = len(folder.FICHIER)
    for subfolder in folder.DOSSIER_:
        total += get_total_file_count(subfolder)
    return total

def get_total_file_count_by_id(folder_id):
    folder = DOSSIER.query.filter_by(id_Dossier=folder_id).first()
    return get_total_file_count(folder)
