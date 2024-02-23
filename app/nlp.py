import fitz
from spacy.lang.fr.stop_words import STOP_WORDS as fr_stop
from spacy.lang.en.stop_words import STOP_WORDS as en_stop
from spacy.util import compile_infix_regex
from spacy import load
from unidecode import unidecode
from docx import Document
from pptx import Presentation
from pandas import read_excel
from concurrent.futures import ThreadPoolExecutor
from collections import Counter
from chardet import detect

def process_pdf(pdf_file):
    """Permet de traiter un fichier PDF et de retourner une liste de mots clés.

    Args:
        pdf_file (str): Le chemin vers le fichier PDF.

    Returns:
        list: La liste des mots clés.
    """
    with fitz.open(pdf_file) as doc:
        text = ''
        for page in doc:
            text += page.get_text()
    return process_nlp(text)

def process_docx(docx_file):
    """Permet de traiter un fichier DOCX et de retourner une liste de mots clés.

    Args:
        docx_file (str): Le chemin vers le fichier DOCX.

    Returns:
        list: La liste des mots clés.
    """
    text = ''
    for paragraph in Document(docx_file).paragraphs:
        text += paragraph.text
    return process_nlp(text)

def process_sheet(sheet_file):
    """Permet de traiter un fichier XLSX et de retourner une liste de mots clés.

    Args:
        sheet_file (str): Le chemin vers le fichier XLSX.

    Returns:
        list: La liste des mots clés.
    """
    text = ''
    text = read_excel(sheet_file).to_string()
    return process_nlp(text)

def process_presentation(presentation_file):
    """Permet de traiter un fichier PPTX et de retourner une liste de mots clés.

    Args:
        presentation_file (str): Le chemin vers le fichier PPTX.

    Returns:
        list: La liste des mots clés.
    """
    text = ''
    for slide in Presentation(presentation_file).slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text
    return process_nlp(text)

def process_txt(txt_file):
    """Permet de traiter un fichier TXT et de retourner une liste de mots clés.

    Args:
        txt_file (str): Le chemin vers le fichier TXT.

    Returns:
        list: La liste des mots clés.
    """
    rawdata = open(txt_file, 'rb').read()
    with open(txt_file, 'r', encoding=detect(rawdata)['encoding']) as file:
        text = file.read()
    return process_nlp(text)

def process_nlp(text, batch_size=100000):
    """Permet de traiter un texte et de retourner une liste de mots clés.

    Args:
        text (str): Le texte à traiter.

    Returns:
        list: La liste des mots clés.
    """
    nlp = load('fr_dep_news_trf')
    nlp.max_length = batch_size
    default_infixes = list(nlp.Defaults.infixes)
    default_infixes.append('[A-Z][a-z0-9]+')
    infix_regex = compile_infix_regex(default_infixes)
    nlp.tokenizer.infix_finditer = infix_regex.finditer
    custom_stop_words = {'document', 'national', 'nationaux', 'france', 'societe', 'societes', 'nom', 'contre', 'voire', 'nan'}

    def clean(text):
        stop_words = fr_stop.union(en_stop).union(custom_stop_words)
        return [unidecode(token.lemma_.lower()) for token in nlp(text)
                if not token.is_stop
                and not token.is_punct
                and not token.is_space
                and not token.like_url
                and not token.like_email
                and len(token) >= 3
                and not token.ent_type_ == 'PER'
                and not token.is_digit
                and not token.is_currency
                and not token.pos_ == 'SYM'
                and not token.pos_ == 'NUM'
                and not token.pos_ == 'VERB'
                and not token.pos_ == 'ADV'
                and not token.lemma_.lower() in stop_words
                and token.lemma_.isalnum()]
    
    def process_batch(batch):
        return Counter(clean(batch))

    batches = [text[i:i + batch_size] for i in range(0, len(text), batch_size)]
    word_frequencies = {}
    with ThreadPoolExecutor() as executor:
        for batch in executor.map(process_batch, batches):
            word_frequencies.update(batch)
    word_frequencies = sorted(word_frequencies.items(), key=lambda x: x[1], reverse=True)
    return [(word, frequency) for word, frequency in word_frequencies]
